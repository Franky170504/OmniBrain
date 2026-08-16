from __future__ import annotations

import csv
import hashlib
import json
import mimetypes
import re
import uuid

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# Constants
# ============================================================

MAX_PAGE_LABEL_LENGTH = 50


SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF",

    ".docx": "DOCX",

    ".pptx": "PPTX",

    ".xlsx": "XLSX",
    ".xlsm": "XLSX",
    ".csv": "XLSX",

    ".txt": "TEXT",
    ".log": "TEXT",
    ".json": "TEXT",
    ".xml": "TEXT",
    ".yaml": "TEXT",
    ".yml": "TEXT",

    ".md": "MARKDOWN",
    ".markdown": "MARKDOWN",

    ".html": "HTML",
    ".htm": "HTML",

    ".png": "IMAGE",
    ".jpg": "IMAGE",
    ".jpeg": "IMAGE",
    ".webp": "IMAGE",
    ".bmp": "IMAGE",
    ".gif": "IMAGE",
    ".tif": "IMAGE",
    ".tiff": "IMAGE",
}

from typing import Any, Iterable, Iterator
import fitz  # PyMuPDF
import pytesseract
from PIL import Image

from app.core.app_config import INPUT_DIR, OUTPUT_DIR

# ============================================================
# Normalized models
# ============================================================

@dataclass
class ParsedPage:
    page_number: int
    text: str

    page_type: str = "DOCUMENT"

    label: str | None = None

    ocr_applied: bool = False


@dataclass
class ParsedChunk:
    chunk_id: str

    document_id: str

    source_file: str

    chunk_index: int

    page_start: int

    page_end: int

    text: str

    character_count: int

    metadata: dict[str, Any] = field(
        default_factory=dict
    )


@dataclass
class ParsedImage:
    image_id: str

    document_id: str

    source_file: str

    page_number: int

    image_index: int

    path: str

    extension: str

    width: int

    height: int

    sha256: str

    caption: str | None = None

    alt_text: str | None = None

    ocr_text: str | None = None

    bbox_x: float = 0.0

    bbox_y: float = 0.0

    bbox_width: float = 1.0

    bbox_height: float = 1.0


@dataclass
class ParsedTable:
    table_id: str

    document_id: str

    source_file: str

    page_number: int

    table_index: int

    title: str | None

    headers: list[str]

    rows: list[list[Any]]

    summary: str | None = None

    storage_format: str = "JSON"

    extraction_engine: str = "omnibrain"

    extraction_confidence: float | None = None

    bbox_x: float = 0.0

    bbox_y: float = 0.0

    bbox_width: float = 1.0

    bbox_height: float = 1.0


@dataclass
class ParsedDocumentResult:
    document: dict[str, Any]

    pages: list[ParsedPage]

    chunks: list[ParsedChunk]

    images: list[ParsedImage]

    tables: list[ParsedTable]


# ============================================================
# Filename / label helpers
# ============================================================

def normalize_display_filename(
    filename: str | None,
) -> str:
    """
    Return a safe original filename for UI/database labels.

    This removes directories but does NOT add temporary
    upload UUID prefixes.
    """

    clean_name = Path(
        filename or "document"
    ).name.strip()

    if not clean_name:
        return "document"

    return clean_name


def filename_page_label(
    filename: str | None,
) -> str:
    """
    Create a page label using the original filename.

    knowledge.pages.page_label is VARCHAR(50).

    Examples:

        report.pdf
            -> report.pdf

        extremely_long_report_name_.....pdf
            -> shortened filename while keeping .pdf
    """

    clean_name = normalize_display_filename(
        filename
    )

    if len(clean_name) <= MAX_PAGE_LABEL_LENGTH:
        return clean_name

    path = Path(
        clean_name
    )

    suffix = path.suffix

    stem = path.stem

    # Keep extension visible.
    if suffix:
        reserved = (
            len(suffix)
            + 3
        )

        available = (
            MAX_PAGE_LABEL_LENGTH
            - reserved
        )
def extract_page_text(document: fitz.Document) -> list[PageText]:
    """
    Extract text from each PDF page.

    Normal PDFs use PyMuPDF text extraction. If a page contains
    no extractable text, render that page as an image and use
    Tesseract OCR as a fallback.
    """
    pages: list[PageText] = []

    for page_index, page in enumerate(document):
        page_number = page_index + 1

        # Fast path: extract embedded PDF text.
        raw_text = page.get_text("text", sort=True)
        normalized_text = normalize_text(raw_text)

        if normalized_text.strip():
            pages.append(
                PageText(
                    page_number=page_number,
                    text=normalized_text,
                )
            )
            continue

        # OCR fallback for scanned/image-only pages.
        LOGGER.info(
            "No extractable text on page %s; running OCR.",
            page_number,
        )

        try:
            # Render page at 2x resolution for better OCR accuracy.
            matrix = fitz.Matrix(2.0, 2.0)
            pixmap = page.get_pixmap(
                matrix=matrix,
                alpha=False,
            )

            image = Image.frombytes(
                "RGB",
                [pixmap.width, pixmap.height],
                pixmap.samples,
            )

            ocr_text = pytesseract.image_to_string(
                image,
                lang="eng",
                config="--psm 3",
            )

            normalized_ocr_text = normalize_text(ocr_text)

            LOGGER.info(
                "OCR completed for page %s: %d characters.",
                page_number,
                len(normalized_ocr_text),
            )

            pages.append(
                PageText(
                    page_number=page_number,
                    text=normalized_ocr_text,
                )
            )

        except Exception as exc:
            LOGGER.exception(
                "OCR failed for page %s: %s",
                page_number,
                exc,
            )

            # Preserve the page in the pipeline even if OCR fails.
            pages.append(
                PageText(
                    page_number=page_number,
                    text="",
                )
            )

    return pages

        if available <= 0:
            return clean_name[
                :MAX_PAGE_LABEL_LENGTH
            ]

        return (
            stem[:available]
            + "..."
            + suffix
        )

    return (
        clean_name[
            :MAX_PAGE_LABEL_LENGTH - 3
        ]
        + "..."
    )


# ============================================================
# Hash / ID helpers
# ============================================================

def sha256_bytes(
    data: bytes,
) -> str:
    return hashlib.sha256(
        data
    ).hexdigest()


def sha256_file(
    path: Path,
) -> str:
    digest = hashlib.sha256()

    with path.open(
        "rb"
    ) as handle:
        while True:
            block = handle.read(
                1024 * 1024
            )

            if not block:
                break

            digest.update(
                block
            )

    return digest.hexdigest()


def stable_uuid(
    *parts: object,
) -> str:
    payload = "\x1f".join(
        str(part)
        for part in parts
    )

    return str(
        uuid.uuid5(
            uuid.NAMESPACE_URL,
            payload,
        )
    )


# ============================================================
# Detection
# ============================================================

def detect_document_type(
    path: Path,
) -> str:
    return SUPPORTED_EXTENSIONS.get(
        path.suffix.lower(),
        "OTHER",
    )


def detect_mime_type(
    path: Path,
) -> str:
    mime_type, _ = mimetypes.guess_type(
        path.name
    )

    return (
        mime_type
        or "application/octet-stream"
    )


# ============================================================
# Text helpers
# ============================================================

def normalize_text(
    text: str,
) -> str:
    text = (
        text
        .replace("\x00", "")
        .replace("\u00ad", "")
    )

    text = re.sub(
        r"[ \t]+",
        " ",
        text,
    )

    text = re.sub(
        r"\n{3,}",
        "\n\n",
        text,
    )

    return text.strip()


def build_table_summary(
    *,
    title: str | None,
    headers: list[str],
    rows: list[list[Any]],
) -> str:
    parts: list[str] = []

    if title:
        parts.append(
            f"Table title: {title}."
        )

    if headers:
        parts.append(
            "Columns: "
            + ", ".join(
                str(value)
                for value in headers
            )
            + "."
        )

    parts.append(
        f"Row count: {len(rows)}."
    )

    if rows:
        preview = "\n".join(
            " | ".join(
                str(value)
                for value in row
            )
            for row in rows[:10]
        )

        parts.append(
            "Table preview:\n"
            + preview
        )

    return "\n".join(
        parts
    )


# ============================================================
# Chunking
# ============================================================

def create_chunks(
    *,
    document_id: str,
    pages: list[ParsedPage],
    source_file: str,
    metadata: dict[str, Any],
    chunk_size: int = 2000,
    overlap: int = 250,
) -> list[ParsedChunk]:

    if chunk_size <= 0:
        raise ValueError(
            "chunk_size must be greater than zero."
        )

    if overlap < 0:
        raise ValueError(
            "overlap cannot be negative."
        )

    if overlap >= chunk_size:
        raise ValueError(
            "overlap must be smaller than chunk_size."
        )

    chunks: list[
        ParsedChunk
    ] = []

    current_text = ""

    current_pages: list[int] = []

    def emit() -> None:
        nonlocal current_text
        nonlocal current_pages

        clean_text = (
            current_text.strip()
        )

        if not clean_text:
            return

        index = (
            len(chunks)
            + 1
        )

        chunks.append(
            ParsedChunk(
                chunk_id=stable_uuid(
                    document_id,
                    "chunk",
                    index,
                    clean_text,
                ),

                document_id=(
                    document_id
                ),

                source_file=(
                    source_file
                ),

                chunk_index=index,

                page_start=min(
                    current_pages
                ),

                page_end=max(
                    current_pages
                ),

                text=clean_text,

                character_count=len(
                    clean_text
                ),

                metadata=metadata,
            )
        )

        if overlap:
            current_text = clean_text[
                -overlap:
            ]
        else:
            current_text = ""

        if current_pages:
            current_pages = [
                current_pages[-1]
            ]
        else:
            current_pages = []

    for page in pages:
        page_text = normalize_text(
            page.text or ""
        )

        if not page_text:
            continue

        projected_size = (
            len(current_text)
            + len(page_text)
            + 2
        )

        if (
            current_text
            and projected_size
            > chunk_size
        ):
            emit()

        if current_text:
            current_text += "\n\n"

        current_text += (
            page_text
        )

        current_pages.append(
            page.page_number
        )

    emit()

    return chunks


# ============================================================
# Document metadata
# ============================================================

def build_document_metadata(
    *,
    file_path: Path,
    display_filename: str,
    document_id: str,
    document_type: str,
    pages: list[ParsedPage],
    chunks: list[ParsedChunk],
    images: list[ParsedImage],
    tables: list[ParsedTable],
) -> dict[str, Any]:

    return {
        "document_id": (
            document_id
        ),

        "metadata": {
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),
        },

        "page_count": len(
            pages
        ),

        "chunk_count": len(
            chunks
        ),

        "image_occurrence_count": len(
            images
        ),

        "unique_image_count": len(
            {
                image.sha256
                for image in images
            }
        ),

        "table_count": len(
            tables
        ),

        "sha256": (
            sha256_file(
                file_path
            )
        ),

        "mime_type": (
            detect_mime_type(
                Path(
                    display_filename
                )
            )
        ),

        "document_type": (
            document_type
        ),
    }


# ============================================================
# Main dispatcher
# ============================================================

def parse_document(
    *,
    file_path: Path,
    output_path: Path,
    display_filename: str | None = None,
    chunk_size: int = 2000,
    overlap: int = 250,
) -> ParsedDocumentResult:
    """
    Parse any supported file.

    file_path:
        Internal temporary file path.

    display_filename:
        ORIGINAL filename uploaded by the user.

    Always use display_filename for page_label.
    """

    original_filename = (
        normalize_display_filename(
            display_filename
            or file_path.name
        )
    )

    document_type = (
        detect_document_type(
            Path(
                original_filename
            )
        )
    )

    extension = (
        Path(
            original_filename
        )
        .suffix
        .lower()
    )

    output_path.mkdir(
        parents=True,
        exist_ok=True,
    )

    if document_type == "PDF":
        return parse_pdf_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if document_type == "DOCX":
        return parse_docx_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if document_type == "PPTX":
        return parse_pptx_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if extension == ".csv":
        return parse_csv_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if document_type == "XLSX":
        return parse_xlsx_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if document_type == "IMAGE":
        return parse_image_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
        )

    if document_type == "HTML":
        return parse_html_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    if document_type in {
        "TEXT",
        "MARKDOWN",
    }:
        return parse_text_document(
            file_path=file_path,
            display_filename=(
                original_filename
            ),
            output_path=(
                output_path
            ),
            document_type=(
                document_type
            ),
            chunk_size=(
                chunk_size
            ),
            overlap=overlap,
        )

    return parse_fallback_document(
        file_path=file_path,
        display_filename=(
            original_filename
        ),
        output_path=(
            output_path
        ),
        chunk_size=(
            chunk_size
        ),
        overlap=overlap,
    )


# ============================================================
# PDF
# ============================================================

def parse_pdf_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    pages: list[
        ParsedPage
    ] = []

    images: list[
        ParsedImage
    ] = []

    tables: list[
        ParsedTable
    ] = []

    image_dir = (
        output_path
        / "images"
        / document_id
    )

    image_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    pdf = fitz.open(
        file_path
    )

    try:
        for page_number, page in enumerate(
            pdf,
            start=1,
        ):

            page_text = normalize_text(
                page.get_text(
                    "text",
                    sort=True,
                )
            )

            pages.append(
                ParsedPage(
                    page_number=(
                        page_number
                    ),

                    page_type=(
                        "DOCUMENT"
                    ),

                    label=label,

                    text=page_text,
                )
            )

            # ----------------------------------------------
            # Images
            # ----------------------------------------------

            for image_index, info in enumerate(
                page.get_images(
                    full=True
                ),
                start=1,
            ):

                try:
                    xref = int(
                        info[0]
                    )

                    image_data = (
                        pdf.extract_image(
                            xref
                        )
                    )

                except Exception:
                    continue

                binary = (
                    image_data.get(
                        "image"
                    )
                )

                if not binary:
                    continue

                image_hash = (
                    sha256_bytes(
                        binary
                    )
                )

                extension = str(
                    image_data.get(
                        "ext"
                    )
                    or "png"
                ).lower()

                width = max(
                    int(
                        image_data.get(
                            "width"
                        )
                        or 1
                    ),
                    1,
                )

                height = max(
                    int(
                        image_data.get(
                            "height"
                        )
                        or 1
                    ),
                    1,
                )

                image_id = (
                    stable_uuid(
                        document_id,
                        "image",
                        page_number,
                        image_index,
                        image_hash,
                    )
                )

                image_path = (
                    image_dir
                    / (
                        f"{image_id}."
                        f"{extension}"
                    )
                )

                image_path.write_bytes(
                    binary
                )

                images.append(
                    ParsedImage(
                        image_id=image_id,

                        document_id=(
                            document_id
                        ),

                        source_file=(
                            display_filename
                        ),

                        page_number=(
                            page_number
                        ),

                        image_index=(
                            image_index
                        ),

                        path=str(
                            image_path
                        ),

                        extension=(
                            extension
                        ),

                        width=width,

                        height=height,

                        sha256=(
                            image_hash
                        ),

                        caption=(
                            page_text[:2000]
                            if page_text
                            else None
                        ),

                        bbox_width=float(
                            width
                        ),

                        bbox_height=float(
                            height
                        ),
                    )
                )

            # ----------------------------------------------
            # Tables
            # ----------------------------------------------

            try:
                finder = (
                    page.find_tables()
                )

                found_tables = (
                    finder.tables
                    if finder
                    else []
                )

            except Exception:
                found_tables = []

            for table_index, table in enumerate(
                found_tables,
                start=1,
            ):

                try:
                    extracted = (
                        table.extract()
                        or []
                    )

                except Exception:
                    extracted = []

                rows: list[
                    list[Any]
                ] = []

                for row in extracted:
                    rows.append(
                        [
                            (
                                ""
                                if value is None
                                else str(
                                    value
                                ).strip()
                            )
                            for value in row
                        ]
                    )

                if not rows:
                    continue

                headers = [
                    str(value)
                    for value
                    in rows[0]
                ]

                bbox = getattr(
                    table,
                    "bbox",
                    None,
                )

                if bbox:
                    x0, y0, x1, y1 = (
                        bbox
                    )

                    bbox_x = max(
                        float(x0),
                        0.0,
                    )

                    bbox_y = max(
                        float(y0),
                        0.0,
                    )

                    bbox_width = max(
                        float(
                            x1 - x0
                        ),
                        1.0,
                    )

                    bbox_height = max(
                        float(
                            y1 - y0
                        ),
                        1.0,
                    )

                else:
                    bbox_x = 0.0
                    bbox_y = 0.0
                    bbox_width = 1.0
                    bbox_height = 1.0

                table_id = (
                    stable_uuid(
                        document_id,
                        "table",
                        page_number,
                        table_index,
                        json.dumps(
                            rows,
                            default=str,
                        ),
                    )
                )

                summary = (
                    build_table_summary(
                        title=(
                            f"Table "
                            f"{table_index}"
                        ),
                        headers=headers,
                        rows=rows,
                    )
                )

                tables.append(
                    ParsedTable(
                        table_id=table_id,

                        document_id=(
                            document_id
                        ),

                        source_file=(
                            display_filename
                        ),

                        page_number=(
                            page_number
                        ),

                        table_index=(
                            table_index
                        ),

                        title=(
                            f"Table "
                            f"{table_index}"
                        ),

                        headers=headers,

                        rows=rows,

                        summary=summary,

                        extraction_engine=(
                            "PyMuPDF"
                        ),

                        bbox_x=bbox_x,

                        bbox_y=bbox_y,

                        bbox_width=(
                            bbox_width
                        ),

                        bbox_height=(
                            bbox_height
                        ),
                    )
                )

    finally:
        pdf.close()

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type="PDF",

            pages=pages,

            chunks=chunks,

            images=images,

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=images,
        tables=tables,
    )


# ============================================================
# DOCX
# ============================================================

def parse_docx_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    doc = DocxDocument(
        str(file_path)
    )

    text_blocks: list[str] = []

    for paragraph in doc.paragraphs:
        value = (
            paragraph.text.strip()
        )

        if value:
            text_blocks.append(
                value
            )

    tables: list[
        ParsedTable
    ] = []

    for table_index, table in enumerate(
        doc.tables,
        start=1,
    ):

        rows = [
            [
                cell.text.strip()
                for cell in row.cells
            ]
            for row in table.rows
        ]

        if not rows:
            continue

        headers = [
            str(value)
            for value in rows[0]
        ]

        summary = build_table_summary(
            title=(
                f"Table {table_index}"
            ),
            headers=headers,
            rows=rows,
        )

        tables.append(
            ParsedTable(
                table_id=stable_uuid(
                    document_id,
                    "table",
                    table_index,
                ),

                document_id=(
                    document_id
                ),

                source_file=(
                    display_filename
                ),

                page_number=1,

                table_index=(
                    table_index
                ),

                title=(
                    f"Table {table_index}"
                ),

                headers=headers,

                rows=rows,

                summary=summary,

                extraction_engine=(
                    "python-docx"
                ),
            )
        )

        text_blocks.append(
            summary
        )

    images: list[
        ParsedImage
    ] = []

    image_dir = (
        output_path
        / "images"
        / document_id
    )

    image_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    image_index = 0

    for relationship in (
        doc.part.rels.values()
    ):
        if (
            "image"
            not in relationship.reltype
        ):
            continue

        try:
            blob = (
                relationship
                .target_part
                .blob
            )

        except Exception:
            continue

        if not blob:
            continue

        image_index += 1

        image_hash = (
            sha256_bytes(
                blob
            )
        )

        content_type = getattr(
            relationship.target_part,
            "content_type",
            "image/png",
        )

        extension = (
            mimetypes.guess_extension(
                content_type
            )
            or ".png"
        ).lstrip(".")

        image_id = stable_uuid(
            document_id,
            "image",
            image_index,
            image_hash,
        )

        image_path = (
            image_dir
            / (
                f"{image_id}."
                f"{extension}"
            )
        )

        image_path.write_bytes(
            blob
        )

        try:
            with PILImage.open(
                image_path
            ) as image:
                width, height = (
                    image.size
                )

        except Exception:
            width = 1
            height = 1

        images.append(
            ParsedImage(
                image_id=image_id,

                document_id=(
                    document_id
                ),

                source_file=(
                    display_filename
                ),

                page_number=1,

                image_index=(
                    image_index
                ),

                path=str(
                    image_path
                ),

                extension=(
                    extension
                ),

                width=max(
                    int(width),
                    1,
                ),

                height=max(
                    int(height),
                    1,
                ),

                sha256=(
                    image_hash
                ),

                caption=(
                    text_blocks[0][:1000]
                    if text_blocks
                    else None
                ),

                bbox_width=float(
                    max(
                        int(width),
                        1,
                    )
                ),

                bbox_height=float(
                    max(
                        int(height),
                        1,
                    )
                ),
            )
        )

    pages = [
        ParsedPage(
            page_number=1,

            page_type=(
                "DOCUMENT"
            ),

            label=label,

            text="\n\n".join(
                text_blocks
            ),
        )
    ]

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type=(
                "DOCX"
            ),

            pages=pages,

            chunks=chunks,

            images=images,

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=images,
        tables=tables,
    )


# ============================================================
# PPTX
# ============================================================

def parse_pptx_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    presentation = Presentation(
        str(file_path)
    )

    pages: list[
        ParsedPage
    ] = []

    tables: list[
        ParsedTable
    ] = []

    images: list[
        ParsedImage
    ] = []

    image_dir = (
        output_path
        / "images"
        / document_id
    )

    image_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    image_counter = 0

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):

        text_blocks: list[str] = []

        table_counter = 0

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text",
            ):
                value = (
                    shape.text.strip()
                )

                if value:
                    text_blocks.append(
                        value
                    )

            if getattr(
                shape,
                "has_table",
                False,
            ):
                table_counter += 1

                rows = [
                    [
                        cell.text.strip()
                        for cell
                        in row.cells
                    ]
                    for row
                    in shape.table.rows
                ]

                if rows:
                    headers = [
                        str(value)
                        for value
                        in rows[0]
                    ]

                    summary = (
                        build_table_summary(
                            title=(
                                f"Table "
                                f"{table_counter}"
                            ),

                            headers=(
                                headers
                            ),

                            rows=rows,
                        )
                    )

                    tables.append(
                        ParsedTable(
                            table_id=(
                                stable_uuid(
                                    document_id,
                                    "table",
                                    slide_number,
                                    table_counter,
                                )
                            ),

                            document_id=(
                                document_id
                            ),

                            source_file=(
                                display_filename
                            ),

                            page_number=(
                                slide_number
                            ),

                            table_index=(
                                table_counter
                            ),

                            title=(
                                f"Table "
                                f"{table_counter}"
                            ),

                            headers=headers,

                            rows=rows,

                            summary=summary,

                            extraction_engine=(
                                "python-pptx"
                            ),
                        )
                    )

                    text_blocks.append(
                        summary
                    )

            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.PICTURE
            ):

                try:
                    blob = (
                        shape.image.blob
                    )

                    extension = (
                        shape.image.ext
                        or "png"
                    )

                except Exception:
                    continue

                image_counter += 1

                image_hash = (
                    sha256_bytes(
                        blob
                    )
                )

                image_id = stable_uuid(
                    document_id,
                    "image",
                    slide_number,
                    image_counter,
                    image_hash,
                )

                image_path = (
                    image_dir
                    / (
                        f"{image_id}."
                        f"{extension}"
                    )
                )

                image_path.write_bytes(
                    blob
                )

                try:
                    with PILImage.open(
                        image_path
                    ) as image:
                        width, height = (
                            image.size
                        )

                except Exception:
                    width = 1
                    height = 1

                images.append(
                    ParsedImage(
                        image_id=(
                            image_id
                        ),

                        document_id=(
                            document_id
                        ),

                        source_file=(
                            display_filename
                        ),

                        page_number=(
                            slide_number
                        ),

                        image_index=(
                            image_counter
                        ),

                        path=str(
                            image_path
                        ),

                        extension=(
                            extension
                        ),

                        width=max(
                            int(width),
                            1,
                        ),

                        height=max(
                            int(height),
                            1,
                        ),

                        sha256=(
                            image_hash
                        ),

                        caption=(
                            "\n".join(
                                text_blocks
                            )[:2000]
                            or None
                        ),

                        bbox_width=float(
                            max(
                                int(width),
                                1,
                            )
                        ),

                        bbox_height=float(
                            max(
                                int(height),
                                1,
                            )
                        ),
                    )
                )

        pages.append(
            ParsedPage(
                page_number=(
                    slide_number
                ),

                page_type="SLIDE",

                label=label,

                text="\n\n".join(
                    text_blocks
                ),
            )
        )

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type=(
                "PPTX"
            ),

            pages=pages,

            chunks=chunks,

            images=images,

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=images,
        tables=tables,
    )


# ============================================================
# XLSX / XLSM
# ============================================================

def parse_xlsx_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    workbook = load_workbook(
        filename=file_path,
        read_only=True,
        data_only=True,
    )

    pages: list[
        ParsedPage
    ] = []

    tables: list[
        ParsedTable
    ] = []

    try:
        for sheet_number, sheet in enumerate(
            workbook.worksheets,
            start=1,
        ):

            rows = [
                [
                    (
                        ""
                        if value is None
                        else value
                    )
                    for value in row
                ]
                for row in sheet.iter_rows(
                    values_only=True
                )
            ]

            rows = [
                row
                for row in rows
                if any(
                    str(value).strip()
                    for value in row
                )
            ]

            if rows:
                headers = [
                    str(value)
                    for value
                    in rows[0]
                ]

                summary = (
                    build_table_summary(
                        title=(
                            sheet.title
                        ),

                        headers=(
                            headers
                        ),

                        rows=rows,
                    )
                )

                tables.append(
                    ParsedTable(
                        table_id=(
                            stable_uuid(
                                document_id,
                                "table",
                                sheet_number,
                            )
                        ),

                        document_id=(
                            document_id
                        ),

                        source_file=(
                            display_filename
                        ),

                        page_number=(
                            sheet_number
                        ),

                        table_index=1,

                        title=(
                            sheet.title
                        ),

                        headers=headers,

                        rows=rows,

                        summary=summary,

                        extraction_engine=(
                            "openpyxl"
                        ),
                    )
                )

                text = summary

            else:
                text = ""

            pages.append(
                ParsedPage(
                    page_number=(
                        sheet_number
                    ),

                    page_type=(
                        "DOCUMENT"
                    ),

                    # Filename for every sheet/page
                    label=label,

                    text=text,
                )
            )

    finally:
        workbook.close()

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type=(
                "XLSX"
            ),

            pages=pages,

            chunks=chunks,

            images=[],

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=[],
        tables=tables,
    )


# ============================================================
# CSV
# ============================================================

def parse_csv_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    with file_path.open(
        "r",
        encoding="utf-8-sig",
        errors="replace",
        newline="",
    ) as handle:

        rows = list(
            csv.reader(
                handle
            )
        )

    headers = (
        rows[0]
        if rows
        else []
    )

    summary = (
        build_table_summary(
            title=(
                display_filename
            ),

            headers=(
                headers
            ),

            rows=rows,
        )
    )

    pages = [
        ParsedPage(
            page_number=1,

            page_type=(
                "DOCUMENT"
            ),

            # Original filename, safely <= 50 chars
            label=label,

            text=summary,
        )
    ]

    tables: list[
        ParsedTable
    ] = []

    if rows:
        tables.append(
            ParsedTable(
                table_id=stable_uuid(
                    document_id,
                    "table",
                    1,
                ),

                document_id=(
                    document_id
                ),

                source_file=(
                    display_filename
                ),

                page_number=1,

                table_index=1,

                title=(
                    display_filename
                ),

                headers=[
                    str(value)
                    for value in headers
                ],

                rows=rows,

                summary=summary,

                extraction_engine="csv",
            )
        )

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type=(
                "XLSX"
            ),

            pages=pages,

            chunks=chunks,

            images=[],

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=[],
        tables=tables,
    )


# ============================================================
# HTML
# ============================================================

def parse_html_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    raw_html = file_path.read_text(
        encoding="utf-8",
        errors="replace",
    )

    soup = BeautifulSoup(
        raw_html,
        "html.parser",
    )

    tables: list[
        ParsedTable
    ] = []

    for table_index, element in enumerate(
        soup.find_all(
            "table"
        ),
        start=1,
    ):

        rows: list[
            list[Any]
        ] = []

        for tr in element.find_all(
            "tr"
        ):
            cells = tr.find_all(
                [
                    "th",
                    "td",
                ]
            )

            row = [
                cell.get_text(
                    " ",
                    strip=True,
                )
                for cell in cells
            ]

            if row:
                rows.append(
                    row
                )

        if not rows:
            continue

        headers = [
            str(value)
            for value in rows[0]
        ]

        summary = (
            build_table_summary(
                title=(
                    f"HTML Table "
                    f"{table_index}"
                ),

                headers=headers,

                rows=rows,
            )
        )

        tables.append(
            ParsedTable(
                table_id=stable_uuid(
                    document_id,
                    "table",
                    table_index,
                ),

                document_id=(
                    document_id
                ),

                source_file=(
                    display_filename
                ),

                page_number=1,

                table_index=(
                    table_index
                ),

                title=(
                    f"HTML Table "
                    f"{table_index}"
                ),

                headers=headers,

                rows=rows,

                summary=summary,

                extraction_engine=(
                    "BeautifulSoup"
                ),
            )
        )

    text = soup.get_text(
        "\n",
        strip=True,
    )

    pages = [
        ParsedPage(
            page_number=1,

            page_type="HTML",

            label=label,

            text=normalize_text(
                text
            ),
        )
    ]

    chunks = create_chunks(
        document_id=document_id,

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type="HTML",

            pages=pages,

            chunks=chunks,

            images=[],

            tables=tables,
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=[],
        tables=tables,
    )


# ============================================================
# TEXT / MARKDOWN / JSON / XML / YAML
# ============================================================

def parse_text_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    document_type: str,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    text = file_path.read_text(
        encoding="utf-8",
        errors="replace",
    )

    if (
        document_type
        == "MARKDOWN"
    ):
        page_type = (
            "MARKDOWN"
        )

    else:
        page_type = (
            "DOCUMENT"
        )

    pages = [
        ParsedPage(
            page_number=1,

            page_type=(
                page_type
            ),

            label=label,

            text=normalize_text(
                text
            ),
        )
    ]

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type=(
                document_type
            ),

            pages=pages,

            chunks=chunks,

            images=[],

            tables=[],
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=[],
        tables=[],
    )


# ============================================================
# IMAGE
# ============================================================

def parse_image_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    with PILImage.open(
        file_path
    ) as image:

        width, height = (
            image.size
        )

    image_id = stable_uuid(
        document_id,
        "image",
        digest,
    )

    image_record = ParsedImage(
        image_id=(
            image_id
        ),

        document_id=(
            document_id
        ),

        source_file=(
            display_filename
        ),

        page_number=1,

        image_index=1,

        path=str(
            file_path
        ),

        extension=(
            Path(
                display_filename
            )
            .suffix
            .lower()
            .lstrip(".")
            or "png"
        ),

        width=max(
            int(width),
            1,
        ),

        height=max(
            int(height),
            1,
        ),

        sha256=digest,

        caption=(
            f"Uploaded image: "
            f"{display_filename}"
        ),

        bbox_width=float(
            max(
                int(width),
                1,
            )
        ),

        bbox_height=float(
            max(
                int(height),
                1,
            )
        ),
    )

    pages = [
        ParsedPage(
            page_number=1,

            page_type="IMAGE",

            label=label,

            text="",
        )
    ]

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type="IMAGE",

            pages=pages,

            chunks=[],

            images=[
                image_record
            ],

            tables=[],
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=[],
        images=[
            image_record
        ],
        tables=[],
    )


# ============================================================
# OTHER / fallback
# ============================================================

def parse_fallback_document(
    *,
    file_path: Path,
    display_filename: str,
    output_path: Path,
    chunk_size: int,
    overlap: int,
) -> ParsedDocumentResult:

    del output_path

    digest = sha256_file(
        file_path
    )

    document_id = stable_uuid(
        "document",
        digest,
    )

    label = filename_page_label(
        display_filename
    )

    try:
        text = file_path.read_text(
            encoding="utf-8",
            errors="replace",
        )

    except Exception:
        text = ""

    pages = [
        ParsedPage(
            page_number=1,

            page_type=(
                "DOCUMENT"
            ),

            label=label,

            text=normalize_text(
                text
            ),
        )
    ]

    chunks = create_chunks(
        document_id=(
            document_id
        ),

        pages=pages,

        source_file=(
            display_filename
        ),

        metadata={
            "title": Path(
                display_filename
            ).stem,

            "filename": (
                display_filename
            ),

            "file_sha256": (
                digest
            ),
        },

        chunk_size=(
            chunk_size
        ),

        overlap=overlap,
    )

    document = (
        build_document_metadata(
            file_path=file_path,

            display_filename=(
                display_filename
            ),

            document_id=(
                document_id
            ),

            document_type="OTHER",

            pages=pages,

            chunks=chunks,

            images=[],

            tables=[],
        )
    )

    return ParsedDocumentResult(
        document=document,
        pages=pages,
        chunks=chunks,
        images=[],
        tables=[],
    )